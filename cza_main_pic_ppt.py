__author__ = 'cza'

import sys
# pyqt5
from PyQt5.QtWidgets import QApplication, QMainWindow, QFileDialog
from PyQt5.QtWidgets import QMessageBox
from PyQt5.QtCore import Qt
# 界面
import cza_pic_ppt_ui

from pynput.keyboard import Listener
import logging  # 日志
from PIL import ImageGrab  # 图像处理

from datetime import *
import os

# import ppt相关库
import pptx
from pptx.util import Inches



class cza_main_pic_ppt(QMainWindow, cza_pic_ppt_ui.Ui_MainWindow):
    def __init__(self):
        global pic_num
        pic_num = 0
        QMainWindow.__init__(self)
        cza_pic_ppt_ui.Ui_MainWindow.__init__(self)
        self.setupUi(self)

        # textEdit_2中显示的内容
        self.textEdit_2.append('欢迎使用本软件.....')
        self.textEdit_2.append('程序开始运行...........')

        # button功能
        self.pushButton.clicked.connect(self.on_pic_ppt)
        self.toolButton.clicked.connect(self.on_open_save_dir)
        self.toolButton_2.clicked.connect(self.on_save_ppt)

        # self.lineEdit.setText("请选择")
        # self.lineEdit_2.setText("请选择")
        # self.pushButton_2.clicked.connect(self.on_scr_print)

    def on_pic_ppt(self):
        print("进行pic到ppt转化")
        # print(FullFileName_save)
        pic_dir = self.lineEdit.text()
        FullSave_dir1 = self.lineEdit_2.text()
        print(FullSave_dir1)
        print(pic_num)

        if pic_num == 0:
            QMessageBox.warning(self, "error", "打开的目录下没有.png图片文件，请确认！！！", QMessageBox.Yes)
        elif FullSave_dir1 == '':
            QMessageBox.warning(self, "出错了！", "未选择图片保存目录", QMessageBox.Yes)
        else:
            print('pic_ppt.....')
            ppt_file = pptx.Presentation()

            pic_files = [fn for fn in os.listdir(pic_dir) if fn.endswith('.png')]
            print(pic_files)

            # for fn in sorted(pic_files, key=lambda item:int(item[:item.index('.')])):
            for fn in sorted(pic_files):
                slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[1])

                # 为PPTX文件当前幻灯片中第一个文本框设置文字，本文代码中可忽略
                slide.shapes.placeholders[0].text = fn[:fn.rindex('.')]

                # 导入并为当前幻灯片添加图片，起始位置和尺寸可修改
                print(fn)
                fn1 = pic_dir + '/' + fn
                print(fn1)
                slide.shapes.add_picture(fn1, Inches(0), Inches(0), Inches(10), Inches(7.5))

            ppt_file.save(FullSave_dir1)
            print("ppt文件已经生成，文件为：" + FullSave_dir1)
            self.textEdit_2.append("ppt文件已经生成，文件为：" + FullSave_dir1)
            QMessageBox.information(self, "提醒", "ppt文件已经生成，文件为：" + FullSave_dir1)

            self.textEdit.clear()

    def on_save_ppt(self):
        global FullFileName_save
        FullFileName_save, _ = QFileDialog.getSaveFileName(self, '文件另存为', r'./', 'ppt (*.pptx)')
        self.lineEdit_2.setText(FullFileName_save)

    def on_open_save_dir(self):
        global pic_num
        print('on_open_save_dir')
        FullSave_dir = QFileDialog.getExistingDirectory(self)
        if len(FullSave_dir) > 0:
            self.lineEdit.setText(FullSave_dir)
            print(FullSave_dir)
            self.textEdit_2.append('图片所在目录为：' + FullSave_dir)
            # 遍历打开目录下的png图片，并显示出来
            pic_files = [fn for fn in os.listdir(FullSave_dir) if fn.endswith('.png')]
            i = 0
            pic_num = len(pic_files)
            if len(pic_files) == 0:
                QMessageBox.warning(self, "出错了！", "打开的目录下没有.png图片文件，请确认！！！", QMessageBox.Yes)
            else:

                str_1 = "打开的目录下，共有：" + str(len(pic_files)) + "张图片待处理！！！"
                print(str_1)
                self.textEdit_2.append(str_1)

            while i < len(pic_files):
                self.textEdit.append(pic_files[i])
                i = i + 1

    def on_scr_print(self):
        print("on_scr_print")
        FullSave_dir = self.lineEdit_2.text()
        if FullSave_dir == '':
            QMessageBox.warning(self, "出错啦！", "未选择图片保存目录", QMessageBox.Yes)
        else:
            self.textEdit.append('开始截屏.......')
            with Listener(on_press=self.press, on_release=self.on_release) as listener:
                listener.join()

    def on_release(self, key):
        pass
        # print('esc')

    def press(self, key):
        logging.info(key)
        print(key, type(key))
        # print(key.char)
        print(str(key))

        if str(key) == "Key.print_screen":
            print("打印屏幕")
            self.screen_print()

        if str(key) == "Key.esc":
            print("退出")
            return False

    def add_edit(self, nr):
        self.textEdit.append('ddd')
        self.textEdit.append(nr)

    def screen_print(self):

        im = ImageGrab.grab()
        dt = datetime.now()
        FullSave_dir = self.lineEdit.text()
        filename = FullSave_dir + "/scr_pic_" + dt.strftime('%Y%m%d%H%M%S') + '.png'
        print(filename)
        self.textEdit.append('成功生成图片：' + filename)
        im.save(filename)


#
# def press(key):
#     logging.info(key)
#     print(key,type(key))
#     #print(key.char)
#     print(str(key))
#
#     if str(key)=="Key.print_screen":
#         print("打印屏幕")
#         screen_print()
#
#     if str(key)=="Key.esc":
#         print("退出")
#         listener.stop()
#         #screen_print()
#
# def screen_print1():
#     im = ImageGrab.grab()
#     dt=datetime.now()
#     filename="screen_pic_"+dt.strftime( '%Y%m%d%H%M%S')+'.png'
#     im.save(filename)
#


if __name__ == '__main__':
    QApplication.setAttribute(Qt.AA_EnableHighDpiScaling)
    app = QApplication(sys.argv)
    md = cza_main_pic_ppt()
    md.show()
    sys.exit(app.exec_())
